import os
import yaml
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
title_slide_layout = prs.slide_layouts[5]

directory = "weekly-reports"
for filename in os.listdir(directory):
    if filename.endswith(".yaml") or filename.endswith(".yml"):
        with open(os.path.join(directory, filename), 'r') as file:
            data = yaml.safe_load(file)

        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = data.get('project', 'Unnamed Project')

        content = f"Status: {data.get('status', '')}\n\nSummary: {data.get('summary', '')}\n\nNext Steps: {data.get('next_steps', '')}\n\nRisks: {data.get('risks', '')}"

        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = content

prs.save("status-report-output.pptx")
